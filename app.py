import os
import re
import json
import base64
import hashlib
from dataclasses import dataclass, asdict
from typing import List, Tuple, Dict, Optional

import numpy as np
import streamlit as st
from dotenv import load_dotenv, find_dotenv

# -------------------- Constants & Paths --------------------
APP_TITLE = "Local RAG POC — Text + Chart Images"
PROJECT_DIR = os.path.dirname(__file__)
DOCS_DIR = os.path.join(PROJECT_DIR, "docs")
CACHE_DIR = os.path.join(PROJECT_DIR, ".rag_cache")
IMG_CACHE = os.path.join(CACHE_DIR, "page_images")
os.makedirs(CACHE_DIR, exist_ok=True)
os.makedirs(IMG_CACHE, exist_ok=True)

# -------------------- Env --------------------
load_dotenv(find_dotenv(), override=False)

# -------------------- Optional deps --------------------
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from openai import OpenAI as _OpenAIClient
except Exception:
    _OpenAIClient = None

try:
    import cohere as _cohere
except Exception:
    _cohere = None

# OCR fallback (optional; requires system Tesseract for best results)
try:
    import pytesseract
    from PIL import Image as _PILImage
except Exception:
    pytesseract = None
    _PILImage = None

# -------------------- Data model --------------------
@dataclass
class Chunk:
    chunk_id: str
    doc_id: str
    filepath: str
    filetype: str          # 'pdf' | 'pptx'
    page_or_slide: int
    text: str
    start_char: int
    end_char: int
    section: str = ""
    doc_title: str = ""
    modality: str = "text"         # "text" | "image_caption"
    image_path: str = ""           # set for page images

# -------------------- Helpers --------------------
def hash_text(x: str) -> str:
    return hashlib.sha256(x.encode("utf-8")).hexdigest()[:16]

def normalize_ws(s: str) -> str:
    return re.sub(r"\s+", " ", s).strip()

def split_sentences(text: str) -> List[str]:
    parts = re.split(r"(?<=[\.!?])\s+(?=[A-Z0-9(])", text)
    return [p.strip() for p in parts if p.strip()]

def sentence_windows(sentences: List[str], max_chars=900, overlap_chars=180) -> List[str]:
    out, cur = [], ""
    for s in sentences:
        if not cur:
            cur = s
        elif len(cur) + 1 + len(s) <= max_chars:
            cur += " " + s
        else:
            out.append(cur)
            cur = (cur[-overlap_chars:] + " " + s) if overlap_chars and len(cur) > overlap_chars else s
    if cur:
        out.append(cur)
    return out

# -------------------- Clients --------------------
def load_openai_client():
    key = os.getenv("OPENAI_API_KEY")
    if key and _OpenAIClient:
        try:
            return _OpenAIClient(api_key=key)
        except Exception:
            return None
    return None

def load_cohere_client():
    key = os.getenv("COHERE_API_KEY")
    if key and _cohere:
        try:
            return _cohere.ClientV2(api_key=key) if hasattr(_cohere, "ClientV2") else _cohere.Client(key)
        except Exception:
            return None
    return None

# -------------------- Ingestion: PDFs & PPTX --------------------
def ingest_pdf(path: str) -> Tuple[List[Chunk], str]:
    if not fitz:
        raise RuntimeError("PyMuPDF (pymupdf) not installed.")
    doc = fitz.open(path)
    doc_id = hash_text(os.path.abspath(path))
    meta = doc.metadata or {}
    title = meta.get("title") or os.path.basename(path)

    chunks: List[Chunk] = []
    for i in range(len(doc)):
        page = doc[i]
        text = normalize_ws(page.get_text("text") or "")
        if not text:
            continue
        windows = sentence_windows(split_sentences(text))
        cursor = 0
        for w in windows:
            start = text.find(w, cursor)
            if start == -1: start = cursor
            end = start + len(w); cursor = end
            chunks.append(Chunk(
                chunk_id=f"{doc_id}-p{i+1}-{hash_text(w)}",
                doc_id=doc_id, filepath=path, filetype="pdf",
                page_or_slide=i+1, text=w, start_char=start, end_char=end,
                section="", doc_title=title, modality="text"
            ))
    doc.close()
    return chunks, title

def extract_text_from_shape(shape):
    parts = []
    try:
        if hasattr(shape, "text") and shape.has_text_frame:
            parts.append(shape.text)
    except Exception: pass
    try:
        if shape.has_table:
            for row in shape.table.rows:
                parts.append(" | ".join(c.text_frame.text if c.text_frame else "" for c in row.cells))
    except Exception: pass
    return "\n".join([p for p in parts if p and p.strip()])

def ingest_pptx(path: str) -> Tuple[List[Chunk], str]:
    if not Presentation:
        raise RuntimeError("python-pptx not installed.")
    prs = Presentation(path)
    doc_id = hash_text(os.path.abspath(path))
    title = os.path.basename(path)
    chunks: List[Chunk] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        slide_title = ""
        try:
            if hasattr(slide.shapes, "title") and slide.shapes.title:
                slide_title = (slide.shapes.title.text or "").strip()
        except Exception: pass
        for shape in slide.shapes:
            t = extract_text_from_shape(shape)
            if t: texts.append(t)
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                n = slide.notes_slide.notes_text_frame.text
                if n: texts.append(n)
        except Exception: pass
        text = normalize_ws("\n".join(texts))
        if not text:
            continue
        windows = sentence_windows(split_sentences(text))
        cursor = 0
        for w in windows:
            start = text.find(w, cursor)
            if start == -1: start = cursor
            end = start + len(w); cursor = end
            chunks.append(Chunk(
                chunk_id=f"{doc_id}-s{idx}-{hash_text(w)}",
                doc_id=doc_id, filepath=path, filetype="pptx",
                page_or_slide=idx, text=w, start_char=start, end_char=end,
                section=slide_title, doc_title=title, modality="text"
            ))
    return chunks, title

def scan_folder(folder: str) -> List[str]:
    paths = []
    for root, _, files in os.walk(folder):
        for f in files:
            if f.lower().endswith((".pdf", ".pptx")):
                paths.append(os.path.join(root, f))
    return sorted(paths)

# -------------------- Page images (PDF) --------------------
def render_pdf_pages_to_images(pdf_path: str, out_dir: str, dpi: int = 180) -> List[str]:
    if not fitz:
        return []
    os.makedirs(out_dir, exist_ok=True)
    doc = fitz.open(pdf_path)
    paths = []
    h = hash_text(pdf_path)
    for i in range(len(doc)):
        out = os.path.join(out_dir, f"{h}_p{i+1}.png")
        if not os.path.exists(out):
            pix = doc[i].get_pixmap(dpi=dpi)
            pix.save(out)
        paths.append(out)
    doc.close()
    return paths

def image_to_data_uri(path: str) -> str:
    # encode image file as a data URI for OpenAI vision
    mime = "image/png"
    with open(path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("utf-8")
    return f"data:{mime};base64,{b64}"

def caption_images_openai(image_paths: List[str], max_chars=400, model="gpt-4o-mini") -> List[str]:
    """
    Generates concise chart/page captions via OpenAI vision.
    Uses inline data URIs so no uploads are necessary.
    """
    client = load_openai_client()
    caps = []
    if not client:
        return [""] * len(image_paths)
    prompt = (
        "You are describing charts for retrieval. "
        "Write one compact caption naming axes, units, series, trends, and any key numbers. "
        f"Maximum {max_chars} characters. No fluff."
    )
    for p in image_paths:
        try:
            resp = client.responses.create(
                model=model,
                input=[{
                    "role": "user",
                    "content": [
                        {"type": "input_text", "text": prompt},
                        {"type": "input_image", "image_url": image_to_data_uri(p)}
                    ],
                }],
                temperature=0.1,
            )
            caps.append((resp.output_text or "").strip())
        except Exception:
            caps.append("")
    return caps

def caption_images_ocr(image_paths: List[str]) -> List[str]:
    if not (pytesseract and _PILImage):
        return [""] * len(image_paths)
    outs = []
    for p in image_paths:
        try:
            txt = pytesseract.image_to_string(_PILImage.open(p))
            outs.append(normalize_ws(txt))
        except Exception:
            outs.append("")
    return outs

def make_image_caption_chunks(pdf_path: str, page_imgs: List[str], captions: List[str], doc_id: str, doc_title: str) -> List[Chunk]:
    out = []
    for i, (img, cap) in enumerate(zip(page_imgs, captions), start=1):
        if not cap:
            continue
        out.append(Chunk(
            chunk_id=f"{doc_id}-imgp{i}-{hash_text(cap)}",
            doc_id=doc_id,
            filepath=pdf_path,
            filetype="pdf",
            page_or_slide=i,
            text=cap,
            start_char=0,
            end_char=len(cap),
            section="(page image)",
            doc_title=doc_title,
            modality="image_caption",
            image_path=img,
        ))
    return out

# -------------------- Embeddings & Rerank --------------------
class CohereV4Embedder:
    def __init__(self, dim: int = 1024):
        self.co = load_cohere_client()
        if not self.co:
            raise RuntimeError("COHERE_API_KEY missing or cohere client not installed.")
        self.dim = dim

    def encode(self, texts: List[str], input_type: str) -> np.ndarray:
        all_vecs = []
        CHUNK = 96
        for i in range(0, len(texts), CHUNK):
            batch = texts[i:i+CHUNK]
            resp = self.co.embed(
                model="embed-v4.0",
                texts=batch,
                input_type=input_type,            # "search_document" | "search_query"
                output_dimension=self.dim,
                embedding_types=["float"],
            )
            floats = getattr(resp.embeddings, "float", None) or getattr(resp.embeddings, "float_", None)
            vecs = np.array(floats, dtype=np.float32)
            norms = np.linalg.norm(vecs, axis=1, keepdims=True) + 1e-9
            vecs = vecs / norms
            all_vecs.append(vecs)
        return np.vstack(all_vecs).astype(np.float32)

class CohereReranker:
    def __init__(self, preferred_model: Optional[str] = None):
        self.co = load_cohere_client()
        if not self.co:
            raise RuntimeError("COHERE_API_KEY missing for rerank.")
        if preferred_model and preferred_model.lower() != "auto":
            self.models = [preferred_model]
        else:
            # Try these in order; adjust to your account access
            self.models = ["rerank-english-v3.0", "rerank-multilingual-v3.0", "rerank-2"]

    def rerank(self, query: str, docs: List[str]) -> List[int]:
        last_err = None
        for m in self.models:
            try:
                resp = self.co.rerank(model=m, query=query, documents=docs)
                ordered = sorted(resp.results, key=lambda x: x.relevance_score, reverse=True)
                return [int(r.index) for r in ordered]
            except Exception as e:
                last_err = e
                continue
        raise RuntimeError(f"All rerank models failed: {last_err}")

# -------------------- Answer generation --------------------
def quote_then_answer_openai(query: str, contexts: List[Chunk], openai_model: str = "gpt-4o-mini") -> str:
    client = load_openai_client()
    if not client:
        bullets = []
        for c in contexts:
            cite = f"[{os.path.basename(c.filepath)}: p.{c.page_or_slide}]"
            bullets.append(f"• {c.text}\n  {cite}")
        return "OpenAI key not set. Top quotes with references:\n\n" + "\n\n".join(bullets)

    system = (
        "You are a careful assistant. Answer ONLY using the provided context snippets.\n"
        "First, extract verbatim quotes that answer the question; then write a short synthesis.\n"
        "Attach inline citations like [filename: page] right after each quote or claim.\n"
        "If the answer is not covered by the context, say you cannot find it."
    )
    ctx = ""
    for i, c in enumerate(contexts, start=1):
        cite = f"[{os.path.basename(c.filepath)}: p.{c.page_or_slide}]"
        ctx += f"\n<CTX {i} cite='{cite}'>\n{c.text}\n</CTX {i}>\n"

    user = f"Question: {query}\n\nContext:\n{ctx}\n\nFormat your answer as:\nQUOTES:\n- \"...\" [file: page]\n\nANSWER:\n... (with citations)"

    try:
        resp = client.responses.create(
            model=openai_model,
            input=[
                {"role": "system", "content": system},
                {"role": "user", "content": user},
            ],
            temperature=0.2,
        )
        return resp.output_text
    except Exception as e:
        return f"OpenAI error: {e}"

def vision_verify_readings(chunks: List[Chunk], question: str, model="gpt-4o-mini") -> List[str]:
    client = load_openai_client()
    if not client:
        return []
    outputs = []
    prompt = (
        "From this chart or page image, extract the values that answer the user question. "
        "Return a short bullet with precise numbers and units if visible."
    )
    for c in chunks:
        if not c.image_path:
            continue
        try:
            resp = client.responses.create(
                model=model,
                input=[{
                    "role": "user",
                    "content": [
                        {"type": "input_text", "text": f"Question: {question}\n{prompt}"},
                        {"type": "input_image", "image_url": image_to_data_uri(c.image_path)}
                    ],
                }],
                temperature=0.1,
            )
            out = (resp.output_text or "").strip()
            if out:
                outputs.append(f"{out} [{os.path.basename(c.filepath)}: p.{c.page_or_slide}]")
        except Exception:
            pass
    return outputs

# -------------------- Persistence: single NPZ --------------------
def index_base(folder: str) -> str:
    return os.path.join(CACHE_DIR, f"index_{hash_text(os.path.abspath(folder))}")

def save_index(folder: str, chunks: List[Chunk],
               text_vecs: Optional[np.ndarray], cap_vecs: Optional[np.ndarray],
               order_text_idxs: List[int], order_cap_idxs: List[int]):
    base = index_base(folder)
    meta = {
        "chunks": [asdict(c) for c in chunks],
        "order_text_idxs": order_text_idxs,
        "order_cap_idxs": order_cap_idxs,
        "dim": int(text_vecs.shape[1]) if text_vecs is not None and text_vecs.size else 1024,
    }
    blob = json.dumps(meta).encode("utf-8")
    np.savez_compressed(
        base + ".npz",
        text=text_vecs if text_vecs is not None else np.zeros((0, meta["dim"]), dtype=np.float32),
        caps=cap_vecs if cap_vecs is not None else np.zeros((0, meta["dim"]), dtype=np.float32),
        meta=np.frombuffer(blob, dtype=np.uint8)
    )

def load_index(folder: str):
    base = index_base(folder) + ".npz"
    if not os.path.exists(base):
        return None
    z = np.load(base, allow_pickle=False)
    meta = json.loads(bytes(z["meta"]).decode("utf-8"))
    chunks = [Chunk(**d) for d in meta["chunks"]]
    text_vecs = z["text"]
    cap_vecs = z["caps"]
    order_text_idxs = meta.get("order_text_idxs", [])
    order_cap_idxs = meta.get("order_cap_idxs", [])
    return chunks, text_vecs, cap_vecs, order_text_idxs, order_cap_idxs

# -------------------- Retrieval utilities --------------------
def cosine_search(qv: np.ndarray, mat: np.ndarray, topk: int) -> List[int]:
    if mat.size == 0:
        return []
    sims = (mat @ qv.reshape(-1, 1)).ravel()
    return [int(i) for i in np.argsort(sims)[::-1][:topk]]

def rrf(list_of_lists: List[List[int]], k: float = 60.0) -> List[int]:
    scores: Dict[int, float] = {}
    for lst in list_of_lists:
        for rank, idx in enumerate(lst, start=1):
            scores[idx] = scores.get(idx, 0.0) + 1.0 / (k + rank)
    return [i for i, _ in sorted(scores.items(), key=lambda x: x[1], reverse=True)]

def mmr_rerank(
    qvec: np.ndarray,
    candidate_ids: List[int],
    text_vecs: np.ndarray,
    cap_vecs: np.ndarray,
    text_map: Dict[int, int],
    cap_map: Dict[int, int],
    topn: int = 8,
    lambda_mult: float = 0.7,
) -> List[int]:
    """
    Maximal Marginal Relevance (MMR) reranker.
    - qvec: normalized query embedding (1D)
    - candidate_ids: global chunk indices to consider
    - text_vecs/cap_vecs: normalized matrices
    - text_map/cap_map: global-id -> local row index
    Returns a new ordering (subset up to topn).
    """
    def vec_for_global(gidx):
        li = text_map.get(gidx)
        if li is not None:
            return text_vecs[li]
        li = cap_map.get(gidx)
        if li is not None:
            return cap_vecs[li]
        return None

    sims_q = {}
    for gid in candidate_ids:
        v = vec_for_global(gid)
        sims_q[gid] = float(np.dot(v, qvec)) if v is not None else -1.0

    selected: List[int] = []
    remaining = set(candidate_ids)

    while remaining and len(selected) < topn:
        best_gid = None
        best_score = -1e9
        for gid in list(remaining):
            v = vec_for_global(gid)
            if v is None:
                continue
            if selected:
                max_sim_to_sel = max(
                    (float(np.dot(v, vec_for_global(s))) for s in selected if vec_for_global(s) is not None),
                    default=0.0
                )
            else:
                max_sim_to_sel = 0.0
            score = lambda_mult * sims_q[gid] - (1.0 - lambda_mult) * max_sim_to_sel
            if score > best_score:
                best_score = score
                best_gid = gid
        if best_gid is None:
            break
        selected.append(best_gid)
        remaining.remove(best_gid)

    return selected

# -------------------- UI --------------------
st.set_page_config(page_title=APP_TITLE, layout="wide")
st.title(APP_TITLE)

with st.sidebar:
    st.header("Settings")
    folder = DOCS_DIR
    # st.text_input("Folder path", value=folder, disabled=True)

    st.subheader("Models")
    openai_model = st.text_input("OpenAI model for caption/answer", value="gpt-4o-mini")

    st.subheader("Retrieval")
    topk_retrieve = st.slider("Retrieve K (before rerank)", 20, 200, 80, step=10)
    topn_context = st.slider("Context N (after rerank)", 3, 20, 8, step=1)
    reranker_mode = st.selectbox("Reranker", ["local (MMR)", "none", "cohere (API)"], index=0)


    if st.button("Build / Rebuild Index", type="primary"):
        if not os.path.isdir(folder):
            st.error("Folder does not exist.")
        else:
            files = scan_folder(folder)
            if not files:
                st.warning("No PDFs or PPTX found under docs/.")
            else:
                st.info(f"Found {len(files)} files. Ingesting…")
                all_chunks: List[Chunk] = []
                # Ingest text
                for p in files:
                    try:
                        if p.lower().endswith(".pdf"):
                            ch, title = ingest_pdf(p)
                            all_chunks.extend(ch)
                        elif p.lower().endswith(".pptx"):
                            ch, title = ingest_pptx(p)
                            all_chunks.extend(ch)
                    except Exception as e:
                        st.warning(f"Failed to parse {os.path.basename(p)}: {e}")

                # Always add image-caption chunks from PDF pages
                for p in [x for x in files if x.lower().endswith(".pdf")]:
                    try:
                        doc_id = hash_text(os.path.abspath(p))
                        _, title = ingest_pdf(p)  # title only; re-opens PDF
                        page_imgs = render_pdf_pages_to_images(p, IMG_CACHE, dpi=180)
                        captions = caption_images_openai(page_imgs, model=openai_model)
                        if any(not c for c in captions):
                            ocr_caps = caption_images_ocr(page_imgs)
                            captions = [c if c else o for c, o in zip(captions, ocr_caps)]
                        img_chunks = make_image_caption_chunks(p, page_imgs, captions, doc_id, title)
                        all_chunks.extend(img_chunks)
                    except Exception as e:
                        st.warning(f"Failed to caption pages for {os.path.basename(p)}: {e}")

                if not all_chunks:
                    st.error("No text or captions extracted from documents.")
                else:
                    st.success(f"Ingested {len(all_chunks)} chunks. Computing embeddings…")
                    text_chunks = [c for c in all_chunks if c.modality == "text"]
                    cap_chunks  = [c for c in all_chunks if c.modality == "image_caption"]
                    order_text_idxs = [all_chunks.index(c) for c in text_chunks]
                    order_cap_idxs  = [all_chunks.index(c) for c in cap_chunks]

                    text_vecs = np.zeros((0, 1024), dtype=np.float32)
                    cap_vecs  = np.zeros((0, 1024), dtype=np.float32)

                    try:
                        embedder = CohereV4Embedder(dim=1024)
                        if text_chunks:
                            text_vecs = embedder.encode([c.text for c in text_chunks], input_type="search_document")
                        if cap_chunks:
                            cap_vecs = embedder.encode([c.text for c in cap_chunks], input_type="search_document")
                    except Exception as e:
                        st.error(f"Embedding failed: {e}")
                        text_vecs = np.zeros((0, 1024), dtype=np.float32)
                        cap_vecs = np.zeros((0, 1024), dtype=np.float32)

                    save_index(folder, all_chunks, text_vecs, cap_vecs, order_text_idxs, order_cap_idxs)
                    st.success("Index saved.")

st.divider()

query = st.text_input("Ask a question about your documents")
ask = st.button("Ask")

if ask and query.strip():
    loaded = load_index(DOCS_DIR)
    if not loaded:
        st.error("No index found. Click 'Build / Rebuild Index' first.")
        st.stop()

    chunks, text_vecs, cap_vecs, order_text_idxs, order_cap_idxs = loaded
    if text_vecs is None:
        st.error("No embeddings matrix found. Rebuild the index.")
        st.stop()

    # Query embedding
    try:
        embedder = CohereV4Embedder(dim=int(text_vecs.shape[1]) if text_vecs.size else 1024)
        qv = embedder.encode([query], input_type="search_query")[0]
    except Exception as e:
        st.error(f"Embedding query failed: {e}")
        st.stop()

    # Retrieve separately for text and captions, then fuse
    text_hits_local = cosine_search(qv, text_vecs, topk=topk_retrieve) if text_vecs.size else []
    cap_hits_local  = cosine_search(qv, cap_vecs,  topk=topk_retrieve) if cap_vecs.size else []

    # Map local indices back to global chunk indices
    text_hit_ids = [order_text_idxs[i] for i in text_hits_local]
    cap_hit_ids  = [order_cap_idxs[i]  for i in cap_hits_local]

    fused = rrf([text_hit_ids, cap_hit_ids])
    candidates = fused[:max(topk_retrieve, topn_context)]

    # Build global->local maps once (for MMR)
    text_map = {g: i for i, g in enumerate(order_text_idxs)}
    cap_map  = {g: i for i, g in enumerate(order_cap_idxs)}

    # Rerank
    if candidates:
        if reranker_mode.startswith("local"):
            order = mmr_rerank(
                qvec=qv,
                candidate_ids=candidates,
                text_vecs=text_vecs,
                cap_vecs=cap_vecs,
                text_map=text_map,
                cap_map=cap_map,
                topn=topn_context,
                lambda_mult=0.7,
            )
            tail = [i for i in candidates if i not in order]
            candidates = order + tail

        elif reranker_mode.startswith("cohere"):
            try:
                rr = CohereReranker(preferred_model="rerank-english-v3.0")
                docs_for_rr = [chunks[i].text for i in candidates]
                order_idx = rr.rerank(query, docs_for_rr)
                candidates = [candidates[i] for i in order_idx]
            except Exception as e:
                st.warning(f"Reranker unavailable: {e}. Falling back to local MMR.")
                order = mmr_rerank(
                    qvec=qv,
                    candidate_ids=candidates,
                    text_vecs=text_vecs,
                    cap_vecs=cap_vecs,
                    text_map=text_map,
                    cap_map=cap_map,
                    topn=topn_context,
                    lambda_mult=0.7,
                )
                tail = [i for i in candidates if i not in order]
                candidates = order + tail
        else:
            # none: keep fused order
            pass

    chosen = candidates[:topn_context]
    chosen_chunks = [chunks[i] for i in chosen]

    # Vision verify is always ON (uses OpenAI; no-op if key missing)
    vision_quotes = vision_verify_readings([c for c in chosen_chunks if c.image_path], query, model=openai_model)

    # Generate answer
    answer = quote_then_answer_openai(query, chosen_chunks, openai_model=openai_model)

    st.subheader("Answer")
    st.write(answer)

    if vision_quotes:
        st.markdown("**Vision reads (chart values):**")
        for ql in vision_quotes:
            st.write(f"- {ql}")

    with st.expander("Sources (top context chunks)"):
        for c in chosen_chunks:
            head = f"{os.path.basename(c.filepath)} — page/slide {c.page_or_slide}"
            if c.section:
                head += f" — {c.section}"
            st.markdown(f"**{head}**")
            if c.image_path:
                st.image(c.image_path, use_container_width=True)
            st.write(c.text)
else:
    st.caption("Build the index, then ask a question here.")
